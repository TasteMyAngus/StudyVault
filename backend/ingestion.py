import os
import sqlite3
from pathlib import Path
from typing import List, Tuple
import hashlib
import uuid
from datetime import datetime
import io
import base64

import pypdf
import docx
import markdown
from .config import DB_PATH, UPLOADS_DIR, CHUNK_SIZE, CHUNK_OVERLAP_PERCENT, OPENAI_API_KEY
from openai import OpenAI


class DocumentParser:
    @staticmethod
    def parse_pdf(file_path: str) -> List[Tuple[int, str]]:
        pages = []
        with open(file_path, 'rb') as f:
            reader = pypdf.PdfReader(f)
            for page_num, page in enumerate(reader.pages, 1):
                text = page.extract_text()
                pages.append((page_num, text))
        return pages
    
    @staticmethod
    def parse_docx(file_path: str) -> List[Tuple[int, str]]:
        # DOCX files are treated as a single page for now.
        doc = docx.Document(file_path)
        text = "\n".join([p.text for p in doc.paragraphs])
        return [(1, text)]
    
    @staticmethod
    def parse_text_file(file_path: str) -> List[Tuple[int, str]]:
        with open(file_path, 'r', encoding='utf-8') as f:
            text = f.read()
        return [(1, text)]
    
    @staticmethod
    def parse_pptx(file_path: str) -> List[Tuple[int, str]]:
        # Fall back to vision when a slide depends on images more than text.
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("Please install python-pptx: pip install python-pptx")
        
        openai_client = OpenAI(api_key=OPENAI_API_KEY)
        pages = []
        
        presentation = Presentation(file_path)
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_text_parts = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text_parts.append(shape.text.strip())
            
            image_descriptions = []
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture type
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        
                        image_b64 = base64.b64encode(image_bytes).decode('utf-8')
                        
                        response = openai_client.chat.completions.create(
                            model="gpt-4-vision-preview",
                            messages=[
                                {
                                    "role": "user",
                                    "content": [
                                        {
                                            "type": "image_url",
                                            "image_url": {
                                                "url": f"data:image/jpeg;base64,{image_b64}"
                                            }
                                        },
                                        {
                                            "type": "text",
                                            "text": "Analyze this diagram or image in detail. Describe what you see, including any labels, relationships, processes, or key information shown. Be comprehensive and descriptive."
                                        }
                                    ]
                                }
                            ],
                            max_tokens=300
                        )
                        
                        image_desc = response.choices[0].message.content
                        image_descriptions.append(f"[IMAGE: {image_desc}]")
                        
                    except Exception as e:
                        print(f"Warning: Could not analyze image on slide {slide_num}: {e}")
                        image_descriptions.append("[IMAGE: Unable to analyze]")
            
            combined_text = "\n".join(slide_text_parts)
            if image_descriptions:
                combined_text += "\n" + "\n".join(image_descriptions)
            
            if combined_text.strip():
                pages.append((slide_num, combined_text))
        
        return pages
    
    @staticmethod
    def parse(file_path: str, doc_type: str) -> List[Tuple[int, str]]:
        if doc_type == "pdf":
            return DocumentParser.parse_pdf(file_path)
        elif doc_type == "docx":
            return DocumentParser.parse_docx(file_path)
        elif doc_type == "pptx":
            return DocumentParser.parse_pptx(file_path)
        elif doc_type in ["txt", "md"]:
            return DocumentParser.parse_text_file(file_path)
        else:
            raise ValueError(f"Unsupported document type: {doc_type}")


class Chunker:
    @staticmethod
    def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap_percent: float = CHUNK_OVERLAP_PERCENT) -> List[str]:
        """Simple word-based chunking (approximation of tokens)."""
        words = text.split()
        overlap = int(chunk_size * overlap_percent)
        chunks = []
        
        for i in range(0, len(words), chunk_size - overlap):
            chunk = " ".join(words[i:i + chunk_size])
            if chunk.strip():
                chunks.append(chunk)
        
        return chunks
    
    @staticmethod
    def chunk_pages(pages: List[Tuple[int, str]]) -> List[dict]:
        import uuid
        chunks = []
        
        for page_num, text in pages:
            page_chunks = Chunker.chunk_text(text)
            for chunk_text in page_chunks:
                chunks.append({
                    "chunk_id": str(uuid.uuid4()),
                    "text": chunk_text,
                    "page_start": page_num,
                    "page_end": page_num,
                    "char_start": None,
                    "char_end": None
                })
        
        return chunks


class Embedder:
    def __init__(self, model_name: str = "all-mpnet-base-v2"):
        try:
            from sentence_transformers import SentenceTransformer
            self.model = SentenceTransformer(model_name)
        except ImportError:
            raise ImportError("Please install sentence-transformers: pip install sentence-transformers")
    
    def embed(self, text: str) -> List[float]:
        embedding = self.model.encode(text)
        return embedding.tolist()
    
    def embed_batch(self, texts: List[str]) -> List[List[float]]:
        embeddings = self.model.encode(texts)
        return embeddings.tolist()


class VectorStore:
    def __init__(self, dimension: int = 768, collection_name: str = "studyvault"):
        try:
            import faiss
            import numpy as np
            self.faiss = faiss
            self.np = np
            self.dimension = dimension
            self.index = faiss.IndexFlatL2(dimension)
            self.id_map = {}
            self.metadata_map = {}
            self.next_id = 0
        except ImportError:
            raise ImportError("Please install faiss-cpu: pip install faiss-cpu")
    
    def add(self, chunk_id: str, embedding: list, metadata: dict, text: str):
        embedding_array = self.np.array([embedding]).astype('float32')
        self.index.add(embedding_array)
        
        current_index = self.next_id
        self.id_map[current_index] = chunk_id
        self.metadata_map[chunk_id] = {
            "text": text,
            "metadata": metadata
        }
        self.next_id += 1
        return current_index
    
    def query(self, embedding: list, top_k: int = 8) -> list:
        embedding_array = self.np.array([embedding]).astype('float32')
        distances, indices = self.index.search(embedding_array, min(top_k, self.next_id))
        
        retrieved = []
        for i, idx in enumerate(indices[0]):
            if idx < 0 or idx >= self.next_id:
                continue
            
            chunk_id = self.id_map[idx]
            stored = self.metadata_map[chunk_id]
            
            retrieved.append({
                "chunk_id": chunk_id,
                "text": stored["text"],
                "score": 1.0 / (1.0 + distances[0][i]),
                "metadata": stored["metadata"]
            })
        
        return retrieved


class IngestionPipeline:
    def __init__(self):
        self.embedder = Embedder()
        self.vector_store = VectorStore()
        self.parser = DocumentParser()
        self.chunker = Chunker()
    
    def ingest(self, file_path: str, doc_id: str, course_id: str, doc_type: str, title: str) -> dict:
        print(f"Parsing {file_path}...")
        pages = self.parser.parse(file_path, doc_type)
        
        print(f"Chunking into ~{CHUNK_SIZE}-token pieces...")
        chunks = self.chunker.chunk_pages(pages)
        print(f"Created {len(chunks)} chunks")
        
        print("Generating embeddings...")
        texts = [c["text"] for c in chunks]
        embeddings = self.embedder.embed_batch(texts)
        
        print("Storing in vector store and database...")
        conn = sqlite3.connect(str(DB_PATH))
        cursor = conn.cursor()
        
        version_id = str(uuid.uuid4())
        file_hash = self._compute_hash(file_path)
        
        cursor.execute('''
            INSERT INTO document_versions 
            (version_id, doc_id, version_label, file_path, file_hash)
            VALUES (?, ?, ?, ?, ?)
        ''', (version_id, doc_id, "v1", file_path, file_hash))
        
        for i, chunk in enumerate(chunks):
            cursor.execute('''
                INSERT INTO chunks
                (chunk_id, version_id, chunk_index, text, page_start, page_end)
                VALUES (?, ?, ?, ?, ?, ?)
            ''', (chunk["chunk_id"], version_id, i, chunk["text"], 
                  chunk["page_start"], chunk["page_end"]))
            
            faiss_index = self.vector_store.add(
                chunk_id=chunk["chunk_id"],
                embedding=embeddings[i],
                metadata={
                    "doc_id": doc_id,
                    "version_id": version_id,
                    "page": chunk["page_start"],
                    "title": title
                },
                text=chunk["text"]
            )
            
            cursor.execute('''
                INSERT INTO embeddings (chunk_id, embedding_ref)
                VALUES (?, ?)
            ''', (chunk["chunk_id"], str(faiss_index)))
        
        conn.commit()
        conn.close()
        
        print(f"Ingestion complete: {len(chunks)} chunks stored")
        return {
            "version_id": version_id,
            "num_chunks": len(chunks),
            "file_hash": file_hash
        }
    
    @staticmethod
    def _compute_hash(file_path: str) -> str:
        sha256 = hashlib.sha256()
        with open(file_path, 'rb') as f:
            for chunk in iter(lambda: f.read(4096), b''):
                sha256.update(chunk)
        return sha256.hexdigest()


if __name__ == "__main__":
    pipeline = IngestionPipeline()
